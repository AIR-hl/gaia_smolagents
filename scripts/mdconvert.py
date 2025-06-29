# This is copied from Magentic-one's great repo: https://github.com/microsoft/autogen/blob/v0.4.4/python/packages/autogen-magentic-one/src/autogen_magentic_one/markdown_browser/mdconvert.py
# Thanks to Microsoft researchers for open-sourcing this!
# type: ignore
import base64
import copy
import html
import json
import mimetypes
import os
import re
import shutil
import subprocess
import sys
import tempfile
import traceback
import zipfile
from typing import Any
from urllib.parse import parse_qs, quote, unquote, urlparse, urlunparse

import mammoth
import markdownify
import pandas as pd
import pptx
import pymupdf4llm
from typing import Any, Optional, Dict, List
# File-format detection
import puremagic
import pydub
import requests
import speech_recognition as sr
from bs4 import BeautifulSoup
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import SRTFormatter
from readability import Document
from smolagents import Tool
import yt_dlp
import pandas as pd
import pandas as pd

from openpyxl import load_workbook

class _CustomMarkdownify(markdownify.MarkdownConverter):
    """
    A custom version of markdownify's MarkdownConverter. Changes include:

    - Altering the default heading style to use '#', '##', etc.
    - Removing javascript hyperlinks.
    - Truncating images with large data:uri sources.
    - Ensuring URIs are properly escaped, and do not conflict with Markdown syntax
    """

    def __init__(self, **options: Any):
        options["heading_style"] = options.get("heading_style", markdownify.ATX)
        super().__init__(**options)

    def convert_hn(self, n: int, el: Any, text: str, convert_as_inline: bool=False, **kwargs) -> str:
        """Same as usual, but be sure to start with a new line"""
        if not convert_as_inline:
            if not re.search(r"^\n", text):
                return "\n" + super().convert_hn(n, el, text, convert_as_inline)

        return super().convert_hn(n, el, text, convert_as_inline)

    def convert_a(self, el: Any, text: str, convert_as_inline: bool=False, **kwargs):
        """Same as usual converter, but removes Javascript links and escapes URIs."""
        prefix, suffix, text = markdownify.chomp(text)
        if not text:
            return ""
        href = el.get("href")
        title = el.get("title")

        if href:
            try:
                parsed_url = urlparse(href)
                if parsed_url.scheme and parsed_url.scheme.lower() not in ["http", "https", "file"]:
                    return "%s%s%s" % (prefix, text, suffix)
                href = urlunparse(parsed_url._replace(path=quote(unquote(parsed_url.path))))
            except ValueError:
                return "%s%s%s" % (prefix, text, suffix)

        if (
            self.options["autolinks"]
            and text.replace(r"\_", "_") == href
            and not title
            and not self.options["default_title"]
        ):
            return "<%s>" % href
        if self.options["default_title"] and not title:
            title = href
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        return "%s[%s](%s%s)%s" % (prefix, text, href, title_part, suffix) if href else text

    def convert_img(self, el: Any, text: str, convert_as_inline: bool=False, **kwargs) -> str:
        """Same as usual converter, but removes data URIs"""
        alt = el.attrs.get("alt", None) or ""
        src = el.attrs.get("src", None) or ""
        title = el.attrs.get("title", None) or ""
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        if convert_as_inline and el.parent.name not in self.options["keep_inline_images_in"]:
            return alt

        if src.startswith("data:"):
            src = src.split(",")[0] + "..."

        return "![%s](%s%s)" % (alt, src, title_part)

    def convert_soup(self, soup: Any) -> str:
        return super().convert_soup(soup)


class DocumentConverterResult:
    """The result of converting a document to text."""
    def __init__(self, title: str | None = None, text_content: str = "", metadata: dict | None = None):
        self.title: str | None = title
        self.text_content: str = text_content
        # 新增metadata字段
        self.metadata: dict = metadata if metadata is not None else {}

# DocumentConverter 和 PlainTextConverter 保持不变
class DocumentConverter:
    """Abstract superclass of all DocumentConverters."""
    def convert(self, local_path: str, **kwargs: Any) -> None | DocumentConverterResult:
        raise NotImplementedError()

class PlainTextConverter(DocumentConverter):
    """Anything with content type text/plain"""
    def convert(self, local_path: str, **kwargs: Any) -> None | DocumentConverterResult:
        content_type, _ = mimetypes.guess_type("__placeholder" + kwargs.get("file_extension", ""))
        if content_type is None:
            return None
        text_content = ""
        # 推荐使用 try-except-finally 来确保文件关闭
        try:
            with open(local_path, "rt", encoding="utf-8") as fh:
                text_content = fh.read()
        except Exception:
            # 可以加入日志记录
            return None
        return DocumentConverterResult(
            title=None,
            text_content=text_content,
        )


class HtmlConverter(DocumentConverter):
    """Anything with content type text/html"""

    def convert(self, local_path: str, **kwargs: Any) -> None | DocumentConverterResult:
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".html", ".htm"]:
            return None
        
        # 建议直接读取字节，让BeautifulSoup处理编码
        try:
            with open(local_path, "rb") as fh:
                return self._convert(fh.read())
        except FileNotFoundError:
            return None
        except Exception as e:
            print(f"Error reading file {local_path}: {e}") # 替换为日志
            return None

    def _convert(self, html_content: str | bytes) -> None | DocumentConverterResult:
        """Helper function that converts an HTML string/bytes."""

        # 1. 使用 readability-lxml 提取主内容
        try:
            doc = Document(html_content)
            main_content_html = doc.summary()
            title = doc.title()
        except Exception as e:
            # readability解析失败时，可以回退到旧的简单解析逻辑
            return self._fallback_convert(html_content)

        # 2. 从原始HTML中提取元数据 (因为readability可能会丢失部分信息)
        soup = BeautifulSoup(html_content, "lxml") # 使用lxml解析器，更快
        metadata = self._extract_metadata(soup)
        
        main_soup = BeautifulSoup(main_content_html, "lxml")
        
        # （可选）在这里可以进行更细致的清理，比如移除主内容里的广告块等
        for ad_block in main_soup.select('.ad, [id*="ad"]'):
            ad_block.decompose()
        
        webpage_text = _CustomMarkdownify(heading_style="atx").convert_soup(main_soup)

        # 4. 清理空白
        # webpage_text = re.sub(r'\n{3,}', '\n\n', webpage_text).strip()
        
        return DocumentConverterResult(
            title=title, 
            text_content=webpage_text,
            metadata=metadata
        )

    def _extract_metadata(self, soup: BeautifulSoup) -> dict:
        """Extracts metadata from the parsed HTML."""
        metadata = {}
        # 作者
        author_meta = soup.find("meta", attrs={"name": "author"})
        if author_meta and author_meta.get("content"):
            metadata["author"] = author_meta["content"]
        
        # 描述
        desc_meta = soup.find("meta", attrs={"name": "description"})
        if desc_meta and desc_meta.get("content"):
            metadata["description"] = desc_meta["content"]

        # 发布时间 (查找几种常见格式)
        pub_time_meta = soup.find("meta", property="article:published_time")
        if pub_time_meta and pub_time_meta.get("content"):
            metadata["published_time"] = pub_time_meta["content"]
        else:
            time_tag = soup.find("time")
            if time_tag and time_tag.get("datetime"):
                metadata["published_time"] = time_tag["datetime"]
        
        return metadata

    def _fallback_convert(self, html_content: str | bytes) -> None | DocumentConverterResult:
        """A fallback converter that uses the original logic if readability fails."""
        soup = BeautifulSoup(html_content, "html.parser")
        for script in soup(["script", "style"]):
            script.decompose() # 使用 decompose() 更彻底
        
        body_elm = soup.find("body")
        if body_elm:
            webpage_text = _CustomMarkdownify().convert_soup(body_elm)
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        return DocumentConverterResult(
            title=soup.title.string if soup.title else None,
            text_content=webpage_text.strip()
        )

class WikipediaConverter(DocumentConverter):
    """
    专门处理维基百科页面，精确提取内容和标题。
    对于这类结构化网站，专用解析器比通用工具更可靠。
    """
    def convert(self, local_path: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        try:
            with open(local_path, "rt", encoding="utf-8") as fh:
                soup = BeautifulSoup(fh.read(), "lxml") # 使用lxml解析器
        except (IOError, UnicodeDecodeError):
            return None

        # 移除脚本和样式
        for element in soup(["script", "style"]):
            element.extract()

        # 精确查找内容和标题元素
        content_elm = soup.find("div", {"id": "mw-content-text"})
        title_elm = soup.find("span", {"class": "mw-page-title-main"})

        main_title = title_elm.get_text(strip=True) if title_elm else \
                     (soup.title.string if soup.title else None)

        if content_elm:
            # 移除目录、信息框等不需要的元素
            for junk_selector in ['.toc', '.infobox', '.mw-editsection']:
                for junk in content_elm.select(junk_selector):
                    junk.extract()
            
            webpage_text = _CustomMarkdownify().convert_soup(content_elm)
            if main_title:
                 webpage_text = f"# {main_title}\n\n{webpage_text}"
        else:
            # 如果找不到特定内容块，则回退到解析body
            body = soup.find("body")
            webpage_text = _CustomMarkdownify().convert_soup(body) if body else ""

        return DocumentConverterResult(title=main_title, text_content=webpage_text.strip())


class YouTubeConverter(DocumentConverter):
    """
    A YouTube converter that directly fetches metadata and transcripts using APIs,
    without relying on a local HTML file. It produces English-only output.
    """
    
    # yt-dlp options to prevent video download and minimize console output
    YDL_OPTS = {
        'noplaylist': True,
        'quiet': True,
        'no_warnings': True,
        'skip_download': True, # Crucial: only extract info, don't download
    }

    def convert(self, url: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        """
        Converts a YouTube URL to a structured DocumentConverterResult.
        The 'local_path' argument is no longer needed.
        """
        # A simple check to see if the URL is likely a YouTube URL
        if "youtube.com/" not in url and "youtu.be/" not in url:
             return None

        # 1. Use yt-dlp to extract video metadata
        try:
            with yt_dlp.YoutubeDL(self.YDL_OPTS) as ydl:
                # The 'url' is used to fetch info directly
                info_dict = ydl.extract_info(url, download=False)
        except yt_dlp.utils.DownloadError as e:
            # This can happen if the video is private, deleted, or the URL is invalid
            print(f"yt-dlp could not extract info from {url}: {e}") # Recommended to replace with a logger
            return None
        
        # 2. Extract desired information from the info_dict
        video_id = info_dict.get("id")
        title = info_dict.get("title", "YouTube Video")
        description = info_dict.get("description", "")
        
        # Extract richer metadata
        channel = info_dict.get("channel")
        channel_url = info_dict.get("channel_url")
        upload_date_str = info_dict.get("upload_date") # Format: YYYYMMDD
        duration_sec = info_dict.get("duration")
        view_count = info_dict.get("view_count")
        like_count = info_dict.get("like_count")
        tags = info_dict.get("tags", [])

        # 3. Format the extracted metadata into English strings
        stats_parts = []
        if view_count is not None:
            stats_parts.append(f"- **Views:** {view_count:,}")
        if like_count is not None:
            stats_parts.append(f"- **Likes:** {like_count:,}")
        if channel and channel_url:
            stats_parts.append(f"- **Channel:** [{channel}]({channel_url})")
        if upload_date_str:
            upload_date = datetime.datetime.strptime(upload_date_str, "%Y%m%d").strftime("%Y-%m-%d")
            stats_parts.append(f"- **Upload Date:** {upload_date}")
        if duration_sec is not None:
            duration_formatted = self._format_duration(duration_sec)
            stats_parts.append(f"- **Duration:** {duration_formatted}")
        if tags:
            stats_parts.append(f"- **Keywords:** {', '.join(tags)}")
        stats = "\n".join(stats_parts)

        # 4. Get the transcript (this logic remains excellent)
        transcript_text = ""
        if video_id:
            try:
                # You can add language preference logic here if needed
                # For example: transcript = YouTubeTranscriptApi.list_transcripts(video_id).find_transcript(['en', 'es'])
                transcript = YouTubeTranscriptApi.get_transcript(video_id)
                transcript_text = SRTFormatter().format_transcript(transcript)
            except Exception as e:
                # e.g., TranscriptsDisabled, NoTranscriptFound
                transcript_text = f"(Could not retrieve transcript: {e})"
        
        # 5. Assemble all content into a final Markdown string
        webpage_text = (
            f"# {title}\n\n"
            f"### Video Metadata\n{stats}\n\n"
            f"### Description\n{description.strip()}\n\n"
            f"### Transcript\n{transcript_text.strip()}\n"
        )

        return DocumentConverterResult(title=title, text_content=webpage_text.strip())

    def _format_duration(self, seconds: int) -> str:
        """Converts total seconds into a human-readable HH:MM:SS string."""
        td = datetime.timedelta(seconds=seconds)
        mm, ss = divmod(td.seconds, 60)
        hh, mm = divmod(mm, 60)
        if td.days > 0:
            hh += td.days * 24
        return f"{hh:02d}:{mm:02d}:{ss:02d}"

class PdfConverter(DocumentConverter):
    """
    Converts PDFs to Markdown. Most style information is ignored, so the results are essentially plain-text.
    """

    def convert(self, local_path, **kwargs) -> None | DocumentConverterResult:
        # Bail if not a PDF
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pdf":
            return None

        return DocumentConverterResult(
            title=None,
            # text_content=pdfminer.high_level.extract_text(local_path),
            text_content=pymupdf4llm.to_markdown(local_path, table_strategy="lines"), # 改用pymupdf
        )

class DocxConverter(HtmlConverter):
    """
    Converts DOCX files to Markdown. Style information (e.g.m headings) and tables are preserved where possible.
    """

    def convert(self, local_path, **kwargs) -> None | DocumentConverterResult:
        # Bail if not a DOCX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".docx":
            return None

        result = None
        with open(local_path, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file)
            html_content = result.value
            result = self._convert(html_content)

        return result


class XlsxConverter(DocumentConverter):
    """
    Converts XLSX files to HTML with style info: bold, italic, font color, background color per cell.
    """

    def get_cell_style(self, cell):
        styles = []
        # 加粗
        if cell.font and cell.font.bold:
            styles.append('font-weight:bold;')
        # 斜体
        if cell.font and cell.font.italic:
            styles.append('font-style:italic;')
        # 文字颜色
        color = getattr(cell.font, 'color', None)
        if color is not None and getattr(color, 'type', None) == 'rgb':
            rgb = getattr(color, 'rgb', None)
            if isinstance(rgb, str) and len(rgb) >= 6:
                styles.append(f'color:#{rgb[-6:]};')
        # 背景色
        fill = getattr(cell, 'fill', None)
        fgColor = getattr(fill, 'fgColor', None)
        if fgColor is not None and getattr(fgColor, 'type', None) == 'rgb':
            rgb = getattr(fgColor, 'rgb', None)
            if isinstance(rgb, str) and rgb != '00000000' and len(rgb) >= 6:
                styles.append(f'background-color:#{rgb[-6:]};')
        return ''.join(styles)

    def convert(self, local_path, **kwargs) -> None | DocumentConverterResult:
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".xlsx", ".xls"]:
            return None

        wb = load_workbook(local_path, data_only=True)
        final_content = ""
        for sheet in wb.worksheets:
            final_content += f"<h2>{sheet.title}</h2>\n<table border='1'>\n"
            for i, row in enumerate(sheet.iter_rows()):
                final_content += "<tr>"
                for cell in row:
                    tag = "th" if i == 0 else "td"
                    style = self.get_cell_style(cell)
                    value = cell.value if cell.value is not None else ""
                    
                    # 只有当有样式时才添加style属性
                    if style:
                        final_content += f"<{tag} style='{style}'>{value}</{tag}>"
                    else:
                        final_content += f"<{tag}>{value}</{tag}>"
                final_content += "</tr>\n"
            final_content += "</table>\n\n"

        return DocumentConverterResult(
            title=None,
            text_content=final_content.strip(),
        )


class PptxConverter(HtmlConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables and images with alt text.
    """

    def convert(self, local_path, **kwargs) -> None | DocumentConverterResult:
        # Bail if not a PPTX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pptx":
            return None

        md_content = ""

        presentation = pptx.Presentation(local_path)
        slide_num = 0
        for slide in presentation.slides:
            slide_num += 1

            md_content += f"\n\n<!-- Slide number: {slide_num} -->\n"

            title = slide.shapes.title
            for shape in slide.shapes:
                # Pictures
                if self._is_picture(shape):
                    # https://github.com/scanny/python-pptx/pull/512#issuecomment-1713100069
                    alt_text = ""
                    try:
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
                    except Exception:
                        pass

                    # A placeholder name
                    filename = re.sub(r"\W", "", shape.name) + ".jpg"
                    md_content += "\n![" + (alt_text if alt_text else shape.name) + "](" + filename + ")\n"

                # Tables
                if self._is_table(shape):
                    html_table = "<html><body><table>"
                    first_row = True
                    for row in shape.table.rows:
                        html_table += "<tr>"
                        for cell in row.cells:
                            if first_row:
                                html_table += "<th>" + html.escape(cell.text) + "</th>"
                            else:
                                html_table += "<td>" + html.escape(cell.text) + "</td>"
                        html_table += "</tr>"
                        first_row = False
                    html_table += "</table></body></html>"
                    md_content += "\n" + self._convert(html_table).text_content.strip() + "\n"

                # Text areas
                elif shape.has_text_frame:
                    if shape == title:
                        md_content += "# " + shape.text.lstrip() + "\n"
                    else:
                        md_content += shape.text + "\n"

            md_content = md_content.strip()

            if slide.has_notes_slide:
                md_content += "\n\n### Notes:\n"
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame is not None:
                    md_content += notes_frame.text
                md_content = md_content.strip()

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _is_picture(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            return True
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
            if hasattr(shape, "image"):
                return True
        return False

    def _is_table(self, shape):
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
        return False


class MediaConverter(DocumentConverter):
    """
    Abstract class for multi-modal media (e.g., images and audio)
    """

    def _get_metadata(self, local_path):
        exiftool = shutil.which("exiftool")
        if not exiftool:
            return None
        else:
            try:
                result = subprocess.run([exiftool, "-json", local_path], capture_output=True, text=True).stdout
                return json.loads(result)[0]
            except Exception:
                return None


class WavConverter(MediaConverter):
    """
    Converts WAV files to markdown via extraction of metadata (if `exiftool` is installed), and speech transcription (if `speech_recognition` is installed).
    """

    def convert(self, local_path, **kwargs) -> None | DocumentConverterResult:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".wav":
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "Title",
                "Artist",
                "Author",
                "Band",
                "Album",
                "Genre",
                "Track",
                "DateTimeOriginal",
                "CreateDate",
                "Duration",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Transcribe
        try:
            transcript = self._transcribe_audio(local_path)
            md_content += "\n\n### Audio Transcript:\n" + ("[No speech detected]" if transcript == "" else transcript)
        except Exception:
            md_content += f"\n\n### Audio Transcript:\nError. Could not transcribe this audio.{e}"

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _transcribe_audio(self, local_path) -> str:
        speech2text = Tool.from_space(
            "hf-audio/whisper-large-v3-turbo",
            name="speech_transcribe",
            description="transcribe the audio to text",
            token="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"
        )
        return speech2text(local_path)


class Mp3Converter(WavConverter):
    """
    Converts MP3 and M4A files to markdown via extraction of metadata (if `exiftool` is installed), and speech transcription (if `speech_recognition` AND `pydub` are installed).
    """

    def convert(self, local_path, **kwargs) -> None | DocumentConverterResult:
        # Bail if not a MP3
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".mp3", ".m4a"]:
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "Title",
                "Artist",
                "Author",
                "Band",
                "Album",
                "Genre",
                "Track",
                "DateTimeOriginal",
                "CreateDate",
                "Duration",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        try:
            transcript = super()._transcribe_audio(local_path).strip()
            md_content += "\n\n### Audio Transcript:\n" + (
                "[No speech detected]" if transcript == "" else transcript
            )
        except Exception as e:
            md_content += f"\n\n### Audio Transcript:\nError. Could not transcribe this audio.{e}"

        # Return the result
        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class ZipConverter(DocumentConverter):
    """
    Extracts ZIP files to a permanent local directory and returns a listing of extracted files.
    """

    def __init__(self, extract_dir: str = "downloads"):
        """
        Initialize with path to extraction directory.

        Args:
            extract_dir: The directory where files will be extracted. Defaults to "downloads"
        """
        self.extract_dir = extract_dir
        # Create the extraction directory if it doesn't exist
        os.makedirs(self.extract_dir, exist_ok=True)

    def convert(self, local_path: str, **kwargs: Any) -> None | DocumentConverterResult:
        # Bail if not a ZIP file
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".zip":
            return None

        # Verify it's actually a ZIP file
        if not zipfile.is_zipfile(local_path):
            return None

        # Extract all files and build list
        extracted_files = []
        with zipfile.ZipFile(local_path, "r") as zip_ref:
            # Extract all files
            zip_ref.extractall(self.extract_dir)
            # Get list of all files
            for file_path in zip_ref.namelist():
                # Skip directories
                if not file_path.endswith("/"):
                    extracted_files.append(self.extract_dir + "/" + file_path)

        # Sort files for consistent output
        extracted_files.sort()

        # Build the markdown content
        md_content = "Downloaded the following files:\n"
        for file in extracted_files:
            md_content += f"* {file}\n"

        return DocumentConverterResult(title="Extracted Files", text_content=md_content.strip())


class ImageConverter(MediaConverter):
    """
    Converts images to markdown via extraction of metadata (if `exiftool` is installed), OCR (if `easyocr` is installed), and description via a multimodal LLM (if an mlm_client is configured).
    """

    def convert(self, local_path, **kwargs) -> None | DocumentConverterResult:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".jpg", ".jpeg", ".png"]:
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(local_path)
        if metadata:
            for f in [
                "ImageSize",
                "Title",
                "Caption",
                "Description",
                "Keywords",
                "Artist",
                "Author",
                "DateTimeOriginal",
                "CreateDate",
                "GPSPosition",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        img_type = "data:image/jpeg;base64," if local_path[-4:] == ".jpg" else "data:image/png;base64,"
        with open(local_path, "rb") as f:
            img_base64 = base64.b64encode(f.read()).decode("utf-8")
        
        question = kwargs.get('question', None)
        if question:
            prompt = f"""Analyze the image carefully and provide appropriate output according to the following question: {question}

            # Content Analysis Rules
            **Tables**: Represent table content and styles as well-structured HTML codes.
            **Geometric Shapes**: Generate vector graphic code (SVG).
            **Complex Graphics**: Provide extremely detailed description.
            **General Images**: Generate comprehensive caption.
            **Math**: Represent formulas in latex code.

            # Output Format
            ## Image Caption
            [Provide the caption of detiald caption of the image]

            ## Answer
            [Provide a answer to the question based on the image content and analysis]"""
        else:
            prompt = """Analyze the image and provide appropriate output based on its content type:

            # Content Analysis Rules

            **Tables**: Extract all content and styling, return as well-structured HTML.
            **Geometric Shapes**: Generate vector graphic code (SVG).
            **Complex Graphics**: Provide extremely detailed description.
            **General Images**: Generate comprehensive caption.
            **Math**: Represent formulas in latex code.

            # Output Format
            ## Image Caption
            [Provide the caption of detiald caption of the image]

            ## Image Content [Optional]
            [Provide the drawing codes of the images if necessary, e.g. table, geometric]"""
        
        payload = {
            "model": "gpt-4o-0806",
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": prompt,
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"{img_type}{img_base64}"
                            }
                        }
                    ],
                },
            ],
            "max_tokens": 8912,
        }    
        headers = {
            "Content-Type": "application/json",
            "Authorization": "Bearer 64268e2b-188f-4e86-9b2a-8542ba3849c8"
        }
        response = requests.post("http://gpt-proxy.jd.com/v1/chat/completions", headers=headers, json=payload)
        try:
            output = response.json()["choices"][0]["message"]["content"]
        except Exception:
            raise Exception(f"Response format unexpected: {response.json()}")
        return DocumentConverterResult(
            title=None,
            text_content=md_content+"\n# Image Detailed Infomation:\n"+output,
        )


class FileConversionException(Exception):
    pass


class UnsupportedFormatException(Exception):
    pass


class MarkdownConverter:
    """(In preview) An extremely simple text-based document reader, suitable for LLM use.
    This reader will convert common file-types or webpages to Markdown."""

    def __init__(
        self,
        requests_session: requests.Session | None = None,
        mlm_client: Any | None = None,
        mlm_model: Any | None = None,
    ):
        if requests_session is None:
            self._requests_session = requests.Session()
        else:
            self._requests_session = requests_session

        self._mlm_client = mlm_client
        self._mlm_model = mlm_model

        self._page_converters: list[DocumentConverter] = []

        self.register_page_converter(PlainTextConverter())
        self.register_page_converter(HtmlConverter())
        self.register_page_converter(WikipediaConverter())
        self.register_page_converter(YouTubeConverter())
        self.register_page_converter(DocxConverter())
        self.register_page_converter(XlsxConverter())
        self.register_page_converter(PptxConverter())
        self.register_page_converter(WavConverter())
        self.register_page_converter(Mp3Converter())
        self.register_page_converter(ImageConverter())
        self.register_page_converter(ZipConverter())
        self.register_page_converter(PdfConverter())

        # Build extension to converter mapping for optimization
        self._build_extension_mapping()

    def _build_extension_mapping(self):
        """Build a mapping from file extensions to their corresponding converters for optimization."""
        self._extension_mapping = {
            '.pdf': [PdfConverter()],
            '.docx': [DocxConverter()],
            '.xlsx': [XlsxConverter()],
            '.xls': [XlsxConverter()],
            '.pptx': [PptxConverter()],
            '.wav': [WavConverter()],
            '.mp3': [Mp3Converter()],
            '.m4a': [Mp3Converter()],
            '.jpg': [ImageConverter()],
            '.jpeg': [ImageConverter()],
            '.png': [ImageConverter()],
            '.zip': [ZipConverter()],
            '.html': [HtmlConverter(),WikipediaConverter(), YouTubeConverter()],  # Order matters for priority
            '.htm': [HtmlConverter(),WikipediaConverter(), YouTubeConverter()],  # Order matters for priority
        }
        
        # For text files, we need to check content type, so we'll handle them separately
        self._text_converter = PlainTextConverter()

    def convert(
        self, source: str | requests.Response, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: deal with kwargs
        """
        Args:
            - source: can be a string representing a path or url, or a requests.response object
            - extension: specifies the file extension to use when interpreting the file. If None, infer from source (path, uri, content-type, etc.)
        """

        # Local path or url
        if isinstance(source, str):
            if source.startswith("http://") or source.startswith("https://") or source.startswith("file://"):
                return self.convert_url(source, **kwargs)
            else:
                return self.convert_local(source, **kwargs)
        # Request response
        elif isinstance(source, requests.Response):
            return self.convert_response(source, **kwargs)

    def convert_local(self, path: str, **kwargs: Any) -> DocumentConverterResult:  # TODO: deal with kwargs
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Get extension alternatives from the path and puremagic
        base, ext = os.path.splitext(path)
        self._append_ext(extensions, ext)
        self._append_ext(extensions, self._guess_ext_magic(path))

        # Convert
        return self._convert(path, extensions, **kwargs)

    # TODO what should stream's type be?
    def convert_stream(self, stream: Any, **kwargs: Any) -> DocumentConverterResult:  # TODO: deal with kwargs
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Save the file locally to a temporary file. It will be deleted before this method exits
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            # Write to the temporary file
            content = stream.read()
            if isinstance(content, str):
                fh.write(content.encode("utf-8"))
            else:
                fh.write(content)
            fh.close()

            # Use puremagic to check for more extension options
            self._append_ext(extensions, self._guess_ext_magic(temp_path))

            # Convert
            result = self._convert(temp_path, extensions, **kwargs)
        # Clean up
        finally:
            try:
                fh.close()
            except Exception:
                pass
            os.unlink(temp_path)

        return result

    def convert_url(self, url: str, **kwargs: Any) -> DocumentConverterResult:  # TODO: fix kwargs type
        # Send a HTTP request to the URL
        user_agent = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/119.0.0.0 Safari/537.36 Edg/119.0.0.0"
        response = self._requests_session.get(url, stream=True, headers={"User-Agent": user_agent})
        response.raise_for_status()
        return self.convert_response(response, **kwargs)

    def convert_response(
        self, response: requests.Response, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO fix kwargs type
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Guess from the mimetype
        content_type = response.headers.get("content-type", "").split(";")[0]
        self._append_ext(extensions, mimetypes.guess_extension(content_type))

        # Read the content disposition if there is one
        content_disposition = response.headers.get("content-disposition", "")
        m = re.search(r"filename=([^;]+)", content_disposition)
        if m:
            base, ext = os.path.splitext(m.group(1).strip("\"'"))
            self._append_ext(extensions, ext)

        # Read from the extension from the path
        base, ext = os.path.splitext(urlparse(response.url).path)
        self._append_ext(extensions, ext)

        # Save the file locally to a temporary file. It will be deleted before this method exits
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            # Download the file
            for chunk in response.iter_content(chunk_size=512):
                fh.write(chunk)
            fh.close()

            # Use puremagic to check for more extension options
            self._append_ext(extensions, self._guess_ext_magic(temp_path))

            # Convert
            result = self._convert(temp_path, extensions, url=response.url)
        except Exception as e:
            print(f"Error in converting: {e}")

        # Clean up
        finally:
            try:
                fh.close()
            except Exception:
                pass
            os.unlink(temp_path)

        return result

    def _convert(self, local_path: str, extensions: list[str | None], **kwargs) -> DocumentConverterResult:
        """
        Optimized conversion method that uses extension mapping to directly select appropriate converters
        instead of trying all converters for each extension.
        """
        error_trace = ""
        
        # Try extensions in priority order
        for ext in extensions + [None]:  # Try last with no extension
            # Prepare kwargs for this extension
            _kwargs = copy.deepcopy(kwargs)
            
            # Overwrite file_extension appropriately
            if ext is None:
                if "file_extension" in _kwargs:
                    del _kwargs["file_extension"]
            else:
                _kwargs.update({"file_extension": ext})
            converters_to_try = []
            
            if ext is not None and ext.lower() in self._extension_mapping:
                # Use mapped converters for known extensions
                converters_to_try = self._extension_mapping[ext.lower()]
            else:
                # For unknown extensions or None, try text converter first, then fallback to all converters
                converters_to_try = [self._text_converter]
                # If text converter fails, we'll fallback to the original approach below
            
            # Try the mapped converters first
            for converter in converters_to_try:
                try:
                    res = converter.convert(local_path, **_kwargs)
                    if res is not None:
                        # Normalize the content
                        res.text_content = "\n".join([line.rstrip() for line in re.split(r"\r?\n", res.text_content)])
                        res.text_content = re.sub(r"\n{3,}", "\n\n", res.text_content)
                        return res
                except Exception:
                    error_trace = ("\n\n" + traceback.format_exc()).strip()
            
            # If mapped converters failed and this is an unknown extension, fallback to trying all converters
            if ext is None or ext.lower() not in self._extension_mapping:
                for converter in self._page_converters:
                    # Skip converters we already tried
                    if converter in converters_to_try:
                        continue
                        
                    try:
                        res = converter.convert(local_path, **_kwargs)
                        if res is not None:
                            # Normalize the content
                            res.text_content = "\n".join([line.rstrip() for line in re.split(r"\r?\n", res.text_content)])
                            res.text_content = re.sub(r"\n{3,}", "\n\n", res.text_content)
                            return res
                    except Exception:
                        error_trace = ("\n\n" + traceback.format_exc()).strip()

        # If we got this far without success, report any exceptions
        if len(error_trace) > 0:
            raise FileConversionException(
                f"Could not convert '{local_path}' to Markdown. File type was recognized as {extensions}. While converting the file, the following error was encountered:\n\n{error_trace}"
            )

        # Nothing can handle it!
        raise UnsupportedFormatException(
            f"Could not parse the file:'{local_path}'. The formats {extensions} are not supported."
        )

        # If we got this far without success, report any exceptions
        if len(error_trace) > 0:
            raise FileConversionException(
                f"Could not convert '{local_path}' to Markdown. File type was recognized as {extensions}. While converting the file, the following error was encountered:\n\n{error_trace}"
            )

        # Nothing can handle it!
        raise UnsupportedFormatException(
            f"Could not convert '{local_path}' to Markdown. The formats {extensions} are not supported."
        )

    def _append_ext(self, extensions, ext):
        """Append a unique non-None, non-empty extension to a list of extensions."""
        if ext is None:
            return
        ext = ext.strip()
        if ext == "":
            return
        # if ext not in extensions:
        if True:
            extensions.append(ext)

    def _guess_ext_magic(self, path):
        """Use puremagic (a Python implementation of libmagic) to guess a file's extension based on the first few bytes."""
        # Use puremagic to guess
        try:
            guesses = puremagic.magic_file(path)
            if len(guesses) > 0:
                ext = guesses[0].extension.strip()
                if len(ext) > 0:
                    return ext
        except FileNotFoundError:
            pass
        except IsADirectoryError:
            pass
        except PermissionError:
            pass
        return None

    def register_page_converter(self, converter: DocumentConverter) -> None:
        """Register a page text converter."""
        self._page_converters.insert(0, converter)
