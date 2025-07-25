# This is copied from Magentic-one's great repo: https://github.com/microsoft/autogen/blob/v0.4.4/python/packages/autogen-magentic-one/src/autogen_magentic_one/markdown_browser/mdconvert.py
# Thanks to Microsoft researchers for open-sourcing this!
# type: ignore
import base64
import copy
import html
import json
import mimetypes
import os
import re
import shutil
import subprocess
import tempfile
import traceback
import zipfile
from typing import Any,Dict
from urllib.parse import quote, unquote, urlparse, urlunparse
import logging
import mammoth
import markdownify
import pptx
import pymupdf4llm
from typing import Any, Optional
import csv
import yaml
import requests
# File-format detection
import puremagic
import textwrap
from bs4 import BeautifulSoup
from youtube_transcript_api import YouTubeTranscriptApi
from youtube_transcript_api.formatters import SRTFormatter
from readability import Document
from smolagents import Tool
import yt_dlp
from Bio import PDB
from openpyxl import load_workbook
import numpy as np
logging.getLogger(__name__)

class _CustomMarkdownify(markdownify.MarkdownConverter):
    """
    A custom version of markdownify's MarkdownConverter. Changes include:

    - Altering the default heading style to use '#', '##', etc.
    - Removing javascript hyperlinks.
    - Truncating images with large data:uri sources.
    - Ensuring URIs are properly escaped, and do not conflict with Markdown syntax
    """

    def __init__(self, **options: Any):
        options["heading_style"] = options.get("heading_style", markdownify.ATX)
        super().__init__(**options)

    def convert_hn(self, n: int, el: Any, text: str, convert_as_inline: bool=False, **kwargs) -> str:
        """Same as usual, but be sure to start with a new line"""
        if not convert_as_inline:
            if not re.search(r"^\n", text):
                return "\n" + super().convert_hn(n, el, text, convert_as_inline)

        return super().convert_hn(n, el, text, convert_as_inline)

    def convert_a(self, el: Any, text: str, convert_as_inline: bool=False, **kwargs):
        """Same as usual converter, but removes Javascript links and escapes URIs."""
        prefix, suffix, text = markdownify.chomp(text)
        if not text:
            return ""
        href = el.get("href")
        title = el.get("title")

        if href:
            try:
                parsed_url = urlparse(href)
                if parsed_url.scheme and parsed_url.scheme.lower() not in ["http", "https", "file"]:
                    return "%s%s%s" % (prefix, text, suffix)
                href = urlunparse(parsed_url._replace(path=quote(unquote(parsed_url.path))))
            except ValueError:
                return "%s%s%s" % (prefix, text, suffix)

        if (
            self.options["autolinks"]
            and text.replace(r"\_", "_") == href
            and not title
            and not self.options["default_title"]
        ):
            return "<%s>" % href
        if self.options["default_title"] and not title:
            title = href
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        return "%s[%s](%s%s)%s" % (prefix, text, href, title_part, suffix) if href else text

    def convert_img(self, el: Any, text: str, convert_as_inline: bool=False, **kwargs) -> str:
        """Same as usual converter, but removes data URIs"""
        alt = el.attrs.get("alt", None) or ""
        src = el.attrs.get("src", None) or ""
        title = el.attrs.get("title", None) or ""
        title_part = ' "%s"' % title.replace('"', r"\"") if title else ""
        if convert_as_inline and el.parent.name not in self.options["keep_inline_images_in"]:
            return alt

        if src.startswith("data:"):
            src = src.split(",")[0] + "..."

        return "![%s](%s%s)" % (alt, src, title_part)

    def convert_soup(self, soup: Any) -> str:
        return super().convert_soup(soup)


class DocumentConverterResult:
    """The result of converting a document to text."""
    def __init__(self, title: str | None = None, text_content: str = "", metadata: dict | None = None):
        self.title: str | None = title
        self.text_content: str = text_content
        # 新增metadata字段
        self.metadata: dict = metadata if metadata is not None else {}

# DocumentConverter 和 PlainTextConverter 保持不变
class DocumentConverter:
    """Abstract superclass of all DocumentConverters."""
    def convert(self, source_path: str, **kwargs: Any) -> None | DocumentConverterResult:
        raise NotImplementedError()

class PlainTextConverter(DocumentConverter):
    """Anything with content type text/plain"""
    def convert(self, source_path: str, **kwargs: Any) -> None | DocumentConverterResult:
        ext = kwargs.get("file_extension", "").lower()
        text_content = ""
        
        try:
            with open(source_path, "r", encoding="utf-8") as fh:
                if ext in ['.txt','.md','.markdown']:
                    text_content = fh.read()
                elif ext == '.csv':
                    reader = csv.reader(fh)
                    header = next(reader)
                    text_content += '| ' + ' | '.join(header) + ' |\n'
                    text_content += '| ' + ' | '.join(['---'] * len(header)) + ' |\n'
                    for row in reader:
                        text_content += '| ' + ' | '.join(row) + ' |\n'
                elif ext in ['.json', '.jsonld']:
                    data = json.load(fh)
                    text_content = json.dumps(data, indent=2, ensure_ascii=False)
                elif ext == '.jsonl':
                    lines = []
                    for line in fh:
                        data = json.loads(line)
                        lines.append(json.dumps(data, indent=2, ensure_ascii=False))
                    text_content = '\n---\n'.join(lines)
                elif ext in ['.yaml', '.yml']:
                    data = yaml.safe_load(fh)
                    text_content = yaml.dump(data, indent=2, allow_unicode=True)
                else:
                    text_content = fh.read()
            return DocumentConverterResult(
                title=None,
                text_content=text_content,
            )                    
        except Exception as e:
            # You might want to log the error here
            logging.error(f"Error processing {source_path} with PlainTextConverter: {e}")
            try:
                # Fallback to plain text reading if structured parsing fails
                with open(source_path, "rt", encoding="utf-8") as fh:
                    text_content = fh.read()
            except Exception as fallback_e:
                logging.error(f"Fallback plain text reading failed for {source_path}: {fallback_e}")
                return None




# class HtmlConverter(DocumentConverter):
#     """Anything with content type text/html"""

#     def convert(self, source_path: str, **kwargs: Any) -> None | DocumentConverterResult:
#         extension = kwargs.get("file_extension", "")
#         if extension.lower() not in [".html", ".htm"]:
#             return None
        
#         # 建议直接读取字节，让BeautifulSoup处理编码
#         try:
#             with open(source_path, "rb") as fh:
#                 return self._convert(fh.read().decode("utf-8", errors="ignore"))
#         except FileNotFoundError:
#             return None
#         except Exception as e:
#             logging.error(f"Error reading file {source_path}: {e}") # 替换为日志
#             return None

#     def _convert(self, html_content: str) -> None | DocumentConverterResult:
#         """Helper function that converts an HTML string/bytes."""

#         # 1. 使用 readability-lxml 提取主内容
#         try:
#             doc = Document(html_content)
#             main_content_html = doc.summary()
#             title = doc.title()
#         except Exception as e:
#             # readability解析失败时，可以回退到旧的简单解析逻辑
#             return self._fallback_convert(html_content)

#         # 2. 从原始HTML中提取元数据 (因为readability可能会丢失部分信息)
#         soup = BeautifulSoup(html_content, "lxml") # 使用lxml解析器，更快
#         metadata = self._extract_metadata(soup)
        
#         main_soup = BeautifulSoup(main_content_html, "lxml")
        
#         # （可选）在这里可以进行更细致的清理，比如移除主内容里的广告块等
#         for ad_block in main_soup.select('.ad, [id*="ad"]'):
#             ad_block.decompose()
        
#         webpage_text = _CustomMarkdownify(heading_style="atx").convert_soup(main_soup)

#         # 4. 清理空白
#         # webpage_text = re.sub(r'\n{3,}', '\n\n', webpage_text).strip()
        
#         return DocumentConverterResult(
#             title=title, 
#             text_content=webpage_text,
#             metadata=metadata
#         )

#     def _extract_metadata(self, soup: BeautifulSoup) -> dict:
#         """Extracts metadata from the parsed HTML."""
#         metadata = {}
#         # 作者
#         author_meta = soup.find("meta", attrs={"name": "author"})
#         if author_meta and author_meta.get("content"):
#             metadata["author"] = author_meta["content"]
        
#         # 描述
#         desc_meta = soup.find("meta", attrs={"name": "description"})
#         if desc_meta and desc_meta.get("content"):
#             metadata["description"] = desc_meta["content"]

#         # 发布时间 (查找几种常见格式)
#         pub_time_meta = soup.find("meta", property="article:published_time")
#         if pub_time_meta and pub_time_meta.get("content"):
#             metadata["published_time"] = pub_time_meta["content"]
#         else:
#             time_tag = soup.find("time")
#             if time_tag and time_tag.get("datetime"):
#                 metadata["published_time"] = time_tag["datetime"]
        
#         return metadata

#     def _fallback_convert(self, html_content: str | bytes) -> None | DocumentConverterResult:
#         """A fallback converter that uses the original logic if readability fails."""
#         soup = BeautifulSoup(html_content, "html.parser")
#         for script in soup(["script", "style"]):
#             script.decompose() # 使用 decompose() 更彻底
        
#         body_elm = soup.find("body")
#         if body_elm:
#             webpage_text = _CustomMarkdownify().convert_soup(body_elm)
#         else:
#             webpage_text = _CustomMarkdownify().convert_soup(soup)

#         return DocumentConverterResult(
#             title=soup.title.string if soup.title else None,
#             text_content=webpage_text.strip()
#         )



class HtmlConverter(DocumentConverter):
    """
    Converts HTML content from either a local file path or a URL.
    It can handle static HTML files as well as the rendered HTML output from
    server-side technologies like PHP, ASP, JSP by fetching the content from a URL.
    """

    # 扩大了支持的扩展名列表，主要用于本地文件场景
    SUPPORTED_EXTENSIONS = {".html", ".htm", ".php", ".asp", ".jsp", ".aspx"}

    def convert(self, source: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        """
        Converts a source which can be a URL or a local file path.

        Args:
            source: The URL of the webpage or the path to the local file.
            **kwargs: Additional arguments, e.g., file_extension.
        """
        # 1. 判断输入是 URL 还是本地文件
        is_url = bool(urlparse(source).scheme in ["http", "https"])

        html_content = ""
        if is_url:
            html_content = self._fetch_from_url(source)
        else:
            # 对于本地文件，可以保留基于扩展名的检查
            extension = kwargs.get("file_extension") or f".{source.split('.')[-1]}"
            if extension.lower() not in self.SUPPORTED_EXTENSIONS:
                logging.warning(f"Unsupported file extension '{extension}' for {source}. Skipping.")
                return None
            html_content = self._read_from_local_file(source)

        if not html_content:
            return None

        return self._parse_html(html_content)

    def _fetch_from_url(self, url: str) -> Optional[str]:
        """Fetches HTML content from a given URL."""
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        }
        try:
            response = requests.get(url, headers=headers, timeout=15)
            response.raise_for_status()  # 如果状态码不是 2xx，则抛出异常

            # 检查内容类型，确保是HTML
            content_type = response.headers.get("Content-Type", "").lower()
            if "text/html" not in content_type:
                logging.warning(f"URL {url} returned non-HTML content-type: {content_type}")
                return None
            
            # requests 会根据响应头和内容智能解码，比手动decode更可靠
            return response.text

        except requests.exceptions.RequestException as e:
            logging.error(f"Error fetching URL {url}: {e}")
            return None

    def _read_from_local_file(self, file_path: str) -> Optional[str]:
        """Reads content from a local file."""
        try:
            with open(file_path, "rb") as fh:
                # 依然建议让BeautifulSoup处理复杂的编码问题
                # 但这里使用 utf-8 with errors='ignore' 作为通用降级方案
                return fh.read().decode("utf-8", errors="ignore")
        except FileNotFoundError:
            logging.error(f"File not found at {file_path}")
            return None
        except Exception as e:
            logging.error(f"Error reading file {file_path}: {e}")
            return None

    def _parse_html(self, html_content: str) -> Optional[DocumentConverterResult]:
        """
        Helper function that parses an HTML string using readability and BeautifulSoup.
        (This is the refactored version of your original _convert method)
        """
        # 1. 使用 readability-lxml 提取主内容
        try:
            doc = Document(html_content)
            main_content_html = doc.summary()
            title = doc.title()
        except Exception:
            # readability解析失败时，回退到简单的解析逻辑
            return self._fallback_parse(html_content)

        # 2. 从原始HTML中提取元数据 (因为readability可能会丢失部分信息)
        soup = BeautifulSoup(html_content, "lxml")
        metadata = self._extract_metadata(soup)
        
        # 3. 清理主内容并转换为Markdown
        main_soup = BeautifulSoup(main_content_html, "lxml")
        
        # （可选）在这里可以进行更细致的清理
        for ad_block in main_soup.select('.ad, [id*="ad"], .advertisement'):
            ad_block.decompose()
        
        webpage_text = _CustomMarkdownify(heading_style="atx").convert_soup(main_soup)

        # 4. 清理多余的空白
        webpage_text = re.sub(r'\n{3,}', '\n\n', webpage_text).strip()
        
        return DocumentConverterResult(
            title=title,
            text_content=webpage_text,
            metadata=metadata
        )

    def _extract_metadata(self, soup: BeautifulSoup) -> dict:
        """Extracts metadata from the parsed HTML. (No changes from original)"""
        metadata = {}
        # 作者
        author_meta = soup.find("meta", attrs={"name": "author"})
        if author_meta and author_meta.get("content"):
            metadata["author"] = author_meta["content"]
        
        # 描述
        desc_meta = soup.find("meta", attrs={"name": "description"})
        if desc_meta and desc_meta.get("content"):
            metadata["description"] = desc_meta["content"]

        # 发布时间 (查找几种常见格式)
        pub_time_meta = soup.find("meta", property="article:published_time")
        if pub_time_meta and pub_time_meta.get("content"):
            metadata["published_time"] = pub_time_meta["content"]
        else:
            time_tag = soup.find("time")
            if time_tag and time_tag.get("datetime"):
                metadata["published_time"] = time_tag["datetime"]
        
        return metadata

    def _fallback_parse(self, html_content: str) -> Optional[DocumentConverterResult]:
        """A fallback parser if readability fails. (No changes from original)"""
        soup = BeautifulSoup(html_content, "lxml") # 建议在fallback中也使用lxml
        for element in soup(["script", "style", "nav", "footer", "aside"]):
            element.decompose()
        
        body_elm = soup.find("body")
        if body_elm:
            webpage_text = _CustomMarkdownify().convert_soup(body_elm)
        else:
            webpage_text = _CustomMarkdownify().convert_soup(soup)

        return DocumentConverterResult(
            title=soup.title.string if soup.title else "Untitled",
            text_content=webpage_text.strip()
        )

class HtmlConverter(DocumentConverter):
    """
    一个更强大的HTML转换器，不仅能处理静态URL和本地文件，
    还能通过模拟表单提交来抓取动态页面的内容。
    """

    def convert(self, source: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        is_url = bool(urlparse(source).scheme in ["http", "https"])
        if is_url:
            html_content = self._fetch_with_get(source)
            if not html_content:
                return None
            return self._parse_html(html_content)
        else:
            # 本地文件处理逻辑 (和之前一样)
            html_content = self._read_from_local_file(source)
            if not html_content:
                return None
            return self._parse_html(html_content)

    def convert_from_form_submission(self, url: str, form_data: Dict[str, Any], method: str = 'GET') -> Optional[DocumentConverterResult]:
        """
        通过模拟提交表单来获取和解析页面。

        Args:
            url (str): 表单将要提交到的URL (即 <form> 标签的 action 属性).
            form_data (Dict[str, Any]): 一个包含表单数据的字典. key是input/select的'name', value是要提交的值.
            method (str, optional): 提交方法, 'GET' 或 'POST'. 默认为 'GET'.

        Returns:
            Optional[DocumentConverterResult]: 解析结果.
        """
        logging.info(f"Submitting form to {url} with method {method} and data {form_data}")
        html_content = self._fetch_with_form(url, form_data, method)
        if not html_content:
            logging.error("Failed to fetch content after form submission.")
            return None
        
        # 获取到内容后，使用和之前完全相同的解析逻辑
        return self._parse_html(html_content)

    def _fetch_with_get(self, url: str) -> Optional[str]:
        """使用简单的GET请求获取内容。"""
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        }
        try:
            response = requests.get(url, headers=headers, timeout=15)
            response.raise_for_status()
            content_type = response.headers.get("Content-Type", "").lower()
            if "text/html" not in content_type:
                logging.warning(f"URL {url} returned non-HTML content-type: {content_type}")
                return None
            return response.text
        except requests.exceptions.RequestException as e:
            logging.error(f"Error during GET request to {url}: {e}")
            return None

    def _fetch_with_form(self, url: str, form_data: Dict[str, Any], method: str) -> Optional[str]:
        """使用 GET 或 POST 请求提交表单数据并获取内容。"""
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36"
        }
        try:
            if method.upper() == 'GET':
                # 对于GET请求，表单数据放在URL参数中
                response = requests.get(url, params=form_data, headers=headers, timeout=20)
            elif method.upper() == 'POST':
                # 对于POST请求，表单数据放在请求体中
                response = requests.post(url, data=form_data, headers=headers, timeout=20)
            else:
                logging.error(f"Unsupported HTTP method: {method}")
                return None
            
            response.raise_for_status()
            return response.text
        except requests.exceptions.RequestException as e:
            logging.error(f"Error during form submission to {url}: {e}")
            return None
    
    def _parse_html(self, html_content: str) -> Optional[DocumentConverterResult]:
        try:
            doc = Document(html_content)
            main_content_html = doc.summary()
            title = doc.title()
        except Exception:
            return self._fallback_parse(html_content)
        soup = BeautifulSoup(html_content, "lxml")
        metadata = self._extract_metadata(soup)
        main_soup = BeautifulSoup(main_content_html, "lxml")
        for ad_block in main_soup.select('.ad, [id*="ad"], .advertisement'):
            ad_block.decompose()
        webpage_text = _CustomMarkdownify(heading_style="atx").convert_soup(main_soup)
        webpage_text = re.sub(r'\n{3,}', '\n\n', webpage_text).strip()
        return DocumentConverterResult(title=title, text_content=webpage_text, metadata=metadata)

    def _extract_metadata(self, soup: BeautifulSoup) -> dict:
        metadata = {}
        author_meta = soup.find("meta", attrs={"name": "author"})
        if author_meta and author_meta.get("content"): metadata["author"] = author_meta["content"]
        desc_meta = soup.find("meta", attrs={"name": "description"})
        if desc_meta and desc_meta.get("content"): metadata["description"] = desc_meta["content"]
        pub_time_meta = soup.find("meta", property="article:published_time")
        if pub_time_meta and pub_time_meta.get("content"):
            metadata["published_time"] = pub_time_meta["content"]
        else:
            time_tag = soup.find("time")
            if time_tag and time_tag.get("datetime"): metadata["published_time"] = time_tag["datetime"]
        return metadata
    def _fallback_parse(self, html_content: str) -> Optional[DocumentConverterResult]:
        soup = BeautifulSoup(html_content, "lxml")
        for element in soup(["script", "style", "nav", "footer", "aside"]): element.decompose()
        body_elm = soup.find("body")
        webpage_text = _CustomMarkdownify().convert_soup(body_elm if body_elm else soup)
        return DocumentConverterResult(title=soup.title.string if soup.title else "Untitled", text_content=webpage_text.strip())
    def _read_from_local_file(self, file_path: str) -> Optional[str]:
        try:
            with open(file_path, "rb") as fh: return fh.read().decode("utf-8", errors="ignore")
        except FileNotFoundError: return None
        except Exception: return None
class PDBConverter(DocumentConverter):
    def convert(self, source_path, **kwargs):
        try:
            parser = PDB.PDBParser(QUIET=True)
            structure = parser.get_structure("protein", source_path)
            
            atoms = list(structure.get_atoms())
            if len(atoms) < 2:
                return "Error: PDB file contains fewer than two atoms."
            
            atom1, atom2 = atoms[0], atoms[1]
            # 正确的距离计算
            distance = np.linalg.norm(atom1.coord - atom2.coord)
            
            result = f"First atom: {atom1.get_name()} ({atom1.coord})\n" \
                    f"Second atom: {atom2.get_name()} ({atom2.coord})\n" \
                    f"Distance_Å: {distance:.3f} Angstroms (Å)"
            return DocumentConverterResult(title=None, text_content=result)
        
        except Exception as e:
            return f"Error parsing PDB file: {str(e)}"

class WikipediaConverter(DocumentConverter):
    """
    专门处理维基百科页面，精确提取内容和标题。
    对于这类结构化网站，专用解析器比通用工具更可靠。
    """
    def convert(self, source_path: str, **kwargs) -> Optional[DocumentConverterResult]:
        try:
            with open(source_path, "rt", encoding="utf-8") as fh:
                soup = BeautifulSoup(fh.read(), "lxml") # 使用lxml解析器
        except (IOError, UnicodeDecodeError):
            return None

        # 移除脚本和样式
        for element in soup(["script", "style"]):
            element.extract()

        # 精确查找内容和标题元素
        content_elm = soup.find("div", {"id": "mw-content-text"})
        title_elm = soup.find("span", {"class": "mw-page-title-main"})

        main_title = title_elm.get_text(strip=True) if title_elm else \
                     (soup.title.string if soup.title else None)

        if content_elm:
            # 移除目录、信息框等不需要的元素
            for junk_selector in ['.toc', '.infobox', '.mw-editsection']:
                for junk in content_elm.select(junk_selector):
                    junk.extract()
            
            webpage_text = _CustomMarkdownify().convert_soup(content_elm)
            if main_title:
                 webpage_text = f"# {main_title}\n\n{webpage_text}"
        else:
            # 如果找不到特定内容块，则回退到解析body
            body = soup.find("body")
            webpage_text = _CustomMarkdownify().convert_soup(body) if body else ""

        return DocumentConverterResult(title=main_title, text_content=webpage_text.strip())


class YouTubeConverter(DocumentConverter):
    """
    A YouTube converter that directly fetches metadata and transcripts using APIs,
    without relying on a local HTML file. It produces English-only output.
    """
    
    # yt-dlp options to prevent video download and minimize console output
    YDL_OPTS = {
        'noplaylist': True,
        'quiet': True,
        'no_warnings': True,
        'skip_download': True, # Crucial: only extract info, don't download
    }

    def convert(self, source_path: str, **kwargs: Any) -> Optional[DocumentConverterResult]:
        """
        Converts a YouTube URL to a structured DocumentConverterResult.
        The 'source_path' argument is no longer needed.
        """
        # A simple check to see if the URL is likely a YouTube URL
        if "youtube.com/" not in source_path and "youtu.be/" not in source_path:
             return None

        # 1. Use yt-dlp to extract video metadata
        try:
            with yt_dlp.YoutubeDL(self.YDL_OPTS) as ydl:
                # The 'source_path' is used to fetch info directly
                info_dict = ydl.extract_info(source_path, download=False)
        except yt_dlp.utils.DownloadError as e:
            # This can happen if the video is private, deleted, or the URL is invalid
            logging.error(f"yt-dlp could not extract info from {source_path}: {e}") # Recommended to replace with a logger
            return None
        
        # 2. Extract desired information from the info_dict
        video_id = info_dict.get("id")
        title = info_dict.get("title", "YouTube Video")
        description = info_dict.get("description", "")
        
        # Extract richer metadata
        channel = info_dict.get("channel")
        channel_url = info_dict.get("channel_url")
        upload_date_str = info_dict.get("upload_date") # Format: YYYYMMDD
        duration_sec = info_dict.get("duration")
        view_count = info_dict.get("view_count")
        like_count = info_dict.get("like_count")
        tags = info_dict.get("tags", [])

        # 3. Format the extracted metadata into English strings
        stats_parts = []
        if view_count is not None:
            stats_parts.append(f"- **Views:** {view_count:,}")
        if like_count is not None:
            stats_parts.append(f"- **Likes:** {like_count:,}")
        if channel and channel_url:
            stats_parts.append(f"- **Channel:** [{channel}]({channel_url})")
        if upload_date_str:
            upload_date = datetime.datetime.strptime(upload_date_str, "%Y%m%d").strftime("%Y-%m-%d")
            stats_parts.append(f"- **Upload Date:** {upload_date}")
        if duration_sec is not None:
            duration_formatted = self._format_duration(duration_sec)
            stats_parts.append(f"- **Duration:** {duration_formatted}")
        if tags:
            stats_parts.append(f"- **Keywords:** {', '.join(tags)}")
        stats = "\n".join(stats_parts)

        # 4. Get the transcript (this logic remains excellent)
        transcript_text = ""
        if video_id:
            try:
                # You can add language preference logic here if needed
                # For example: transcript = YouTubeTranscriptApi.list_transcripts(video_id).find_transcript(['en', 'es'])
                transcript = YouTubeTranscriptApi.get_transcript(video_id)
                transcript_text = SRTFormatter().format_transcript(transcript)
            except Exception as e:
                # e.g., TranscriptsDisabled, NoTranscriptFound
                transcript_text = f"(Could not retrieve transcript: {e})"
        
        # 5. Assemble all content into a final Markdown string
        webpage_text = (
            f"# {title}\n\n"
            f"### Video Metadata\n{stats}\n\n"
            f"### Description\n{description.strip()}\n\n"
            f"### Transcript\n{transcript_text.strip()}\n"
        )

        return DocumentConverterResult(title=title, text_content=webpage_text.strip())

    def _format_duration(self, seconds: int) -> str:
        """Converts total seconds into a human-readable HH:MM:SS string."""
        td = datetime.timedelta(seconds=seconds)
        mm, ss = divmod(td.seconds, 60)
        hh, mm = divmod(mm, 60)
        if td.days > 0:
            hh += td.days * 24
        return f"{hh:02d}:{mm:02d}:{ss:02d}"

class PdfConverter(DocumentConverter):
    """
    Converts PDFs to Markdown. Most style information is ignored, so the results are essentially plain-text.
    """

    def convert(self, source_path, **kwargs) -> None | DocumentConverterResult:
        # Bail if not a PDF
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".pdf":
            return None

        return DocumentConverterResult(
            title=None,
            # text_content=pdfminer.high_level.extract_text(source_path),
            text_content=pymupdf4llm.to_markdown(source_path, table_strategy="lines"), # 改用pymupdf
        )

class DocxConverter(HtmlConverter):
    """
    Converts DOCX files to Markdown. Style information (e.g.m headings) and tables are preserved where possible.
    """

    def convert(self, source_path, **kwargs) -> None | DocumentConverterResult:
        # Bail if not a DOCX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".docx":
            return None

        result = None
        with open(source_path, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file)
            html_content = result.value
            result = self._convert(html_content)

        return result


class XlsxConverter(DocumentConverter):
    """
    Converts XLSX files to HTML with style info: bold, italic, font color, background color per cell.
    """

    def get_cell_style(self, cell):
        styles = []
        # 加粗
        if cell.font and cell.font.bold:
            styles.append('font-weight:bold;')
        # 斜体
        if cell.font and cell.font.italic:
            styles.append('font-style:italic;')
        # 文字颜色
        color = getattr(cell.font, 'color', None)
        if color is not None and getattr(color, 'type', None) == 'rgb':
            rgb = getattr(color, 'rgb', None)
            if isinstance(rgb, str) and len(rgb) >= 6:
                styles.append(f'color:#{rgb[-6:]};')
        # 背景色
        fill = getattr(cell, 'fill', None)
        fgColor = getattr(fill, 'fgColor', None)
        if fgColor is not None and getattr(fgColor, 'type', None) == 'rgb':
            rgb = getattr(fgColor, 'rgb', None)
            if isinstance(rgb, str) and rgb != '00000000' and len(rgb) >= 6:
                styles.append(f'background-color:#{rgb[-6:]};')
        return ''.join(styles)

    def convert(self, source_path, **kwargs) -> None | DocumentConverterResult:
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".xlsx", ".xls"]:
            return None

        wb = load_workbook(source_path, data_only=True)
        final_content = ""
        for sheet in wb.worksheets:
            final_content += f"<h2>{sheet.title}</h2>\n<table border='1'>\n"
            for i, row in enumerate(sheet.iter_rows()):
                final_content += "<tr>"
                for cell in row:
                    tag = "th" if i == 0 else "td"
                    style = self.get_cell_style(cell)
                    value = cell.value if cell.value is not None else ""
                    
                    # 只有有样式时才添加style属性
                    if style and style!='color:#000000;':
                        final_content += f"<{tag} style='{style}'>{value}</{tag}>"
                    else:
                        final_content += f"<{tag}>{value}</{tag}>"
                final_content += "</tr>\n"
            final_content += "</table>\n\n"

        return DocumentConverterResult(
            title=None,
            text_content=final_content.strip(),
        )


class PptxConverter(HtmlConverter):
    """
    Converts PPTX files to Markdown. Supports heading, tables and images with alt text.
    """
    def _is_picture(self, shape):
        """Check if a shape is a picture."""
        try:
            import pptx
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
                return True
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PLACEHOLDER:
                if hasattr(shape, "image"):
                    return True
        except:
            pass
        return False

    def _is_table(self, shape):
        """Check if a shape is a table."""
        try:
            import pptx
            if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
                return True
        except:
            pass
        return False

    def _extract_image_content(self, shape, slide_num, shape_num):
        """Extract image content and analyze it."""
        try:
            # Try to get alt text first
            alt_text = ""
            try:
                alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")
            except Exception:
                pass

            if alt_text:
                return f"<p><strong>Image {slide_num}-{shape_num}:</strong> {alt_text}</p>"
            
            # If no alt text, try to extract and analyze the image
            try:
                # Extract image data
                image = shape.image
                image_bytes = image.blob
                
                # Save to temporary file
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
                    temp_file.write(image_bytes)
                    temp_image_path = temp_file.name

                # Use parse_image to analyze the image
                from .parse_image import parse_image
                image_analysis = parse_image.invoke(temp_image_path, "Please describe the content of this image in detail")
                
                # Clean up temporary file
                os.unlink(temp_image_path)
                
                return f"<div><strong>Image {slide_num}-{shape_num} Analysis:</strong><br/>{image_analysis}</div>"
                
            except Exception as img_error:
                return f"<p><strong>Image {slide_num}-{shape_num}:</strong> [Cannot parse image content: {str(img_error)}]</p>"
                
        except Exception as e:
            return f"<p><strong>Image {slide_num}-{shape_num}:</strong> [Image processing error: {str(e)}]</p>"


    def convert(self, source_path, **kwargs) -> None | DocumentConverterResult:
        try:
            if not os.path.exists(source_path):
                return f"Error: PowerPoint file '{source_path}' does not exist."
            
            supported_formats = ['.pptx', '.ppt']
            file_ext = os.path.splitext(source_path)[1].lower()
            
            if file_ext not in supported_formats:
                return f"Error: Unsupported file format '{file_ext}'. Supported formats: {', '.join(supported_formats)}"
            
            try:
                import pptx
                
                presentation = pptx.Presentation(source_path)
                html_content = f"<h1>PowerPoint: {os.path.basename(source_path)}</h1>\n"
                html_content += f"<p>Number of slides: {len(presentation.slides)}</p>\n"
                html_content += "<hr/>\n"
                
                slide_num = 0
                for slide in presentation.slides:
                    slide_num += 1
                    html_content += f"<h2>Slide {slide_num}</h2>\n"
                    
                    title = slide.shapes.title
                    shape_num = 0
                    
                    for shape in slide.shapes:
                        shape_num += 1
                        
                        # Process image
                        if self._is_picture(shape):
                            image_content = self._extract_image_content(shape, slide_num, shape_num)
                            html_content += image_content + "\n"
                        
                        # Process table
                        elif self._is_table(shape):
                            html_content += "<h4>Table:</h4>\n"
                            table_content = _extract_table_content(shape)
                            html_content += table_content + "\n"
                        
                        # Process text
                        elif shape.has_text_frame and shape.text.strip():
                            if shape == title:
                                html_content += f"<h3>{html.escape(shape.text.strip())}</h3>\n"
                            else:
                                # Process multi-level text
                                text_content = shape.text.strip()
                                if text_content:
                                    # Split text by line, keep format
                                    lines = text_content.split('\n')
                                    html_content += "<div>\n"
                                    for line in lines:
                                        if line.strip():
                                            html_content += f"<p>{html.escape(line.strip())}</p>\n"
                                    html_content += "</div>\n"
                    
                    # Process slide notes
                    if slide.has_notes_slide:
                        notes_frame = slide.notes_slide.notes_text_frame
                        if notes_frame is not None and notes_frame.text.strip():
                            html_content += "<h4>Notes:</h4>\n"
                            notes_text = notes_frame.text.strip()
                            html_content += f"<div style='background-color:#f5f5f5;padding:10px;'>{html.escape(notes_text)}</div>\n"
                    
                    html_content += "<hr/>\n"
                return DocumentConverterResult(
                    title=None,
                    text_content=html_content.strip(),
                )
            except Exception as e:
                raise Exception(f"Error processing PowerPoint file: {str(e)}")
        except FileNotFoundError:
            raise Exception(f"Error: PowerPoint file '{source_path}' does not exist.")


class MediaConverter(DocumentConverter):
    """
    Abstract class for multi-modal media (e.g., images and audio)
    """

    def _get_metadata(self, source_path):
        exiftool = shutil.which("exiftool")
        if not exiftool:
            return None
        else:
            try:
                result = subprocess.run([exiftool, "-json", source_path], capture_output=True, text=True).stdout
                return json.loads(result)[0]
            except Exception:
                return None


class WavConverter(MediaConverter):
    """
    Converts WAV files to markdown via extraction of metadata (if `exiftool` is installed), and speech transcription (if `speech_recognition` is installed).
    """

    def convert(self, source_path, **kwargs) -> None | DocumentConverterResult:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".wav":
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(source_path)
        if metadata:
            for f in [
                "Title",
                "Artist",
                "Author",
                "Band",
                "Album",
                "Genre",
                "Track",
                "DateTimeOriginal",
                "CreateDate",
                "Duration",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        # Transcribe
        try:
            transcript = self._transcribe_audio(source_path)
            md_content += "\n\n# Audio Transcript:\n" + ("[No speech detected]" if transcript == "" else transcript)
        except Exception:
            md_content += f"\n\n# Error. Could not transcribe this audio. {e}"

        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )

    def _transcribe_audio(self, source_path) -> str:
        speech2text = Tool.from_space(
            "hf-audio/whisper-large-v3-turbo",
            name="speech_transcribe",
            description="transcribe the audio to text",
            token="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"
        )
        return speech2text(source_path)


class Mp3Converter(WavConverter):
    """
    Converts MP3 and M4A files to markdown via extraction of metadata (if `exiftool` is installed), and speech transcription (if `speech_recognition` AND `pydub` are installed).
    """

    def convert(self, source_path, **kwargs) -> None | DocumentConverterResult:
        # Bail if not a MP3
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".mp3", ".m4a"]:
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(source_path)
        if metadata:
            for f in [
                "Title",
                "Artist",
                "Author",
                "Band",
                "Album",
                "Genre",
                "Track",
                "DateTimeOriginal",
                "CreateDate",
                "Duration",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        try:
            transcript = super()._transcribe_audio(source_path).strip()
            md_content += "\n\n### Audio Transcript:\n" + (
                "[No speech detected]" if transcript == "" else transcript
            )
        except Exception as e:
            md_content += f"\n\n### Audio Transcript:\nError. Could not transcribe this audio.{e}"

        # Return the result
        return DocumentConverterResult(
            title=None,
            text_content=md_content.strip(),
        )


class ZipConverter(DocumentConverter):
    """
    Extracts ZIP files to a permanent local directory and returns a listing of extracted files.
    """

    def __init__(self, extract_dir: str = "downloads"):
        """
        Initialize with path to extraction directory.

        Args:
            extract_dir: The directory where files will be extracted. Defaults to "downloads"
        """
        self.extract_dir = extract_dir
        # Create the extraction directory if it doesn't exist
        os.makedirs(self.extract_dir, exist_ok=True)

    def convert(self, source_path: str, **kwargs: Any) -> None | DocumentConverterResult:
        # Bail if not a ZIP file
        extension = kwargs.get("file_extension", "")
        if extension.lower() != ".zip":
            return None

        # Verify it's actually a ZIP file
        if not zipfile.is_zipfile(source_path):
            return None

        # Extract all files and build list
        extracted_files = []
        with zipfile.ZipFile(source_path, "r") as zip_ref:
            # Extract all files
            zip_ref.extractall(self.extract_dir)
            # Get list of all files
            for file_path in zip_ref.namelist():
                # Skip directories
                if not file_path.endswith("/"):
                    extracted_files.append(self.extract_dir + "/" + file_path)

        # Sort files for consistent output
        extracted_files.sort()

        # Build the markdown content
        md_content = "Downloaded the following files:\n"
        for file in extracted_files:
            md_content += f"* {file}\n"

        return DocumentConverterResult(title="Extracted Files", text_content=md_content.strip())


class ImageConverter(MediaConverter):
    """
    Converts images to markdown via extraction of metadata (if `exiftool` is installed), OCR (if `easyocr` is installed), and description via a multimodal LLM (if an mlm_client is configured).
    """

    def convert(self, source_path, **kwargs) -> None | DocumentConverterResult:
        # Bail if not a XLSX
        extension = kwargs.get("file_extension", "")
        if extension.lower() not in [".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp", ".tiff"]:
            return None

        md_content = ""

        # Add metadata
        metadata = self._get_metadata(source_path)
        if metadata:
            for f in [
                "ImageSize",
                "Title",
                "Caption",
                "Description",
                "Keywords",
                "Artist",
                "Author",
                "DateTimeOriginal",
                "CreateDate",
                "GPSPosition",
            ]:
                if f in metadata:
                    md_content += f"{f}: {metadata[f]}\n"

        img_type = "data:image/jpeg;base64," if source_path[-4:] == ".jpg" else "data:image/png;base64,"
        with open(source_path, "rb") as f:
            img_base64 = base64.b64encode(f.read()).decode("utf-8")

        ocr_result = None
        try:
            prompt = textwrap.dedent(f"""
            You are a powerful OCR assistant, please carefully analyze and accurately extract the text, tables or formulas in the picture.\n
            You can use HTML code to represent the tables, LaTex code to represent the formulas. Note: DO NOT add any redundant description.\n
            If there is no such content, only return: "There isn't any text in the picture".
            """)
            payload = {
                "model": "Doubao-1.5-vision-pro-32k",
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": prompt,
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"{img_type}{img_base64}"
                                }
                            }
                        ],
                    },
                ],
                "max_tokens": 4096,
            }    
            headers = {
                "Content-Type": "application/json",
                "Authorization": "Bearer 64268e2b-188f-4e86-9b2a-8542ba3849c8"
            }
            response = requests.post("http://gpt-proxy.jd.com/v1/chat/completions", headers=headers, json=payload)
            if response.json().get("error"):
                raise Exception(response.json()["error"])
            ocr_result = response.json()["choices"][0]["message"]["content"]
        except Exception as e:
            raise RuntimeError(f"Error parsing image with OCR: {str(e)}")

        if ocr_result and "There isn't any text in the picture" in ocr_result:
            ocr_prompt=""
        else:
            ocr_prompt=f"""**OCR Extracted Result:**
This is the OCR result of this image, which is very accurate for the text content. Please refer carefully
```
{ocr_result.strip()}
```
"""

        question = kwargs.get('question', None)
        if question:
            prompt = f"""You are an expert for Visual Question Answering (VQA) task. 
Please analyze the image carefully and provide accurate caption and answer for the following question:
**Question:**
```
{question}
```

**Content Format Rules:**
For different content types, process them according to the following rules:
- Tables: Analyze the content and styling, return as well-structured HTML codes.
- Geometric Shapes: Generate vector graphic code (SVG).
- Complex Graphics: Provide extremely detailed description.
- General Images: Generate detailed and comprehensive caption.
- Math: Represent formulas in LaTex code.

{ocr_prompt}

**Output Format:**
## Image Caption:
[Provide an accurate caption of the image, dont miss any key information.]

## Answer
[Provide a accurate answer to the question based on the image content and above analysis]
"""
        else:
            prompt =f"""You are an expert for analyzing images. 
Please analyze the image carefully and provide accurate caption based on following rules:

**Content Format Rules:**
- Tables: Analyze the content and styling, return as well-structured HTML codes.
- Geometric Shapes: Generate vector graphic code (SVG).
- Complex Graphics: Provide extremely detailed description.
- General Images: Generate detailed and comprehensive caption.
- Math: Represent formulas in LaTex code.

{ocr_prompt}

**Output Format:**
## Image Caption:
[Provide an accurate caption of the image, dont miss any key information]
"""

        payload = {
            "model": "gpt-4o-0806",
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": prompt,
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"{img_type}{img_base64}"
                            }
                        }
                    ],
                },
            ],
            "max_tokens": 16384,
        }    
        headers = {
            "Content-Type": "application/json",
            "Authorization": "Bearer 64268e2b-188f-4e86-9b2a-8542ba3849c8"
        }
        response = requests.post("http://gpt-proxy.jd.com/v1/chat/completions", headers=headers, json=payload)
        try:
            output = response.json()["choices"][0]["message"]["content"]
            output = output.replace("[Submit Button for user to submit their answers.]", "")
            output += "\n\n---\n\n## Additional OCR Result:\n" + ocr_result
        except Exception:
            raise Exception(f"Response format unexpected: {response.json()}")
        return DocumentConverterResult(
            title=None,
            text_content=output,
        )


class FileConversionException(Exception):
    pass


class UnsupportedFormatException(Exception):
    pass


class MarkdownConverter:
    """(In preview) An extremely simple text-based document reader, suitable for LLM use.
    This reader will convert common file-types or webpages to Markdown."""

    def __init__(
        self,
        requests_session: requests.Session | None = None,
        mlm_client: Any | None = None,
        mlm_model: Any | None = None,
    ):
        if requests_session is None:
            self._requests_session = requests.Session()
        else:
            self._requests_session = requests_session

        self._mlm_client = mlm_client
        self._mlm_model = mlm_model

        self._page_converters: list[DocumentConverter] = []

        self.register_page_converter(PlainTextConverter())
        self.register_page_converter(HtmlConverter())
        self.register_page_converter(WikipediaConverter())
        self.register_page_converter(YouTubeConverter())
        self.register_page_converter(DocxConverter())
        self.register_page_converter(XlsxConverter())
        self.register_page_converter(PptxConverter())
        self.register_page_converter(WavConverter())
        self.register_page_converter(Mp3Converter())
        self.register_page_converter(ImageConverter())
        self.register_page_converter(ZipConverter())
        self.register_page_converter(PdfConverter())
        self.register_page_converter(PDBConverter())

        # Build extension to converter mapping for optimization
        self._build_extension_mapping()

    def _build_extension_mapping(self):
        """Build a mapping from file extensions to their corresponding converters for optimization."""
        self._extension_mapping = {
            '.pdf': [PdfConverter()],
            '.docx': [DocxConverter()],
            '.xlsx': [XlsxConverter()],
            '.xls': [XlsxConverter()],
            '.pptx': [PptxConverter()],
            '.wav': [WavConverter()],
            '.mp3': [Mp3Converter()],
            '.m4a': [Mp3Converter()],
            '.jpg': [ImageConverter()],
            '.jpeg': [ImageConverter()],
            '.png': [ImageConverter()],
            '.webp': [ImageConverter()],
            '.gif': [ImageConverter()],
            '.bmp': [ImageConverter()],
            '.tiff': [ImageConverter()],
            '.zip': [ZipConverter()],
            '.txt': [PlainTextConverter()],
            '.srt':[PlainTextConverter()],
            '.html': [HtmlConverter(), YouTubeConverter(), WikipediaConverter()],  # Order matters for priority
            '.htm': [HtmlConverter(), YouTubeConverter(), WikipediaConverter()],  # Order matters for priority
            '.pdb':[PDBConverter()],
        }
        
        # For text files, we need to check content type, so we'll handle them separately
        self._text_converter = PlainTextConverter()

    def convert(
        self, source: str | requests.Response, **kwargs: Any
    ) -> DocumentConverterResult:  # TODO: deal with kwargs
        """
        Args:
            - source: can be a string representing a path or url, or a requests.response object
            - extension: specifies the file extension to use when interpreting the file. If None, infer from source (path, uri, content-type, etc.)
        """

        # Local path or url
        if isinstance(source, str):
            if source.startswith("http://") or source.startswith("https://") or source.startswith("file://"):
                return self.convert_url(source, **kwargs)
            else:
                return self.convert_local(source, **kwargs)
        # Request response
        elif isinstance(source, requests.Response):
            return self.convert_response(source, **kwargs)

    def convert_local(self, path: str, **kwargs: Any) -> DocumentConverterResult:  # TODO: deal with kwargs
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Get extension alternatives from the path and puremagic
        base, ext = os.path.splitext(path)
        self._append_ext(extensions, ext)
        self._append_ext(extensions, self._guess_ext_magic(path))

        # Convert
        return self._convert(path, extensions, **kwargs)

    # TODO what should stream's type be?
    def convert_stream(self, stream: Any, **kwargs: Any) -> DocumentConverterResult:  # TODO: deal with kwargs
        # Prepare a list of extensions to try (in order of priority)
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []

        # Save the file locally to a temporary file. It will be deleted before this method exits
        handle, temp_path = tempfile.mkstemp()
        fh = os.fdopen(handle, "wb")
        result = None
        try:
            # Write to the temporary file
            content = stream.read()
            if isinstance(content, str):
                fh.write(content.encode("utf-8"))
            else:
                fh.write(content)
            fh.close()

            # Use puremagic to check for more extension options
            self._append_ext(extensions, self._guess_ext_magic(temp_path))

            # Convert
            result = self._convert(temp_path, extensions, **kwargs)
        # Clean up
        finally:
            try:
                fh.close()
            except Exception:
                pass
            os.unlink(temp_path)

        return result

    def convert_url(self, url: str, **kwargs: Any) -> DocumentConverterResult:  # TODO: fix kwargs type
        # Send a HTTP request to the URL
        user_agent = "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/137.0.0.0 Safari/537.36 Edg/137.0.0.0"
        response = self._requests_session.get(url, stream=True, headers={"User-Agent": user_agent})
        response.raise_for_status()
        return self.convert_response(response, **kwargs)
    def convert_response(
        self, response: requests.Response, **kwargs: Any
    ) -> DocumentConverterResult:
        # 1. 收集可能的文件扩展名
        ext = kwargs.get("file_extension")
        extensions = [ext] if ext is not None else []
        content_type = response.headers.get("content-type", "").split(";", 1)[0]
        self._append_ext(extensions, mimetypes.guess_extension(content_type))

        dispo = response.headers.get("content-disposition", "")
        m = re.search(r'filename=([^;]+)', dispo)
        if m:
            _, ext = os.path.splitext(m.group(1).strip("\"'"))
            self._append_ext(extensions, ext)

        parsed = urlparse(response.url)
        _, ext = os.path.splitext(parsed.path)
        self._append_ext(extensions, ext)

        # 2. 把响应内容原样写入二进制临时文件
        handle, temp_path = tempfile.mkstemp()
        with os.fdopen(handle, "wb") as fh:
            for chunk in response.iter_content(chunk_size=8192):
                fh.write(chunk)

        try:
            # 3. 再次探测真实扩展名（puremagic）
            # self._append_ext(extensions, self._guess_ext_magic(temp_path))
            # with open(temp_path, "rb") as f:
            #     raw = f.read()
            # text = raw.decode("utf-8", errors="ignore")

            result = self._convert(temp_path, extensions, url=response.url)
        except Exception as e:
            logging.error(f"Error in converting: {e}")
            result = None
        finally:
            os.unlink(temp_path)

        return result

    def _convert(self, source_path: str, extensions: list[str | None], **kwargs) -> DocumentConverterResult:
        """
        Optimized conversion method that uses extension mapping to directly select appropriate converters
        instead of trying all converters for each extension.
        """
        error_trace = ""
        
        # Try extensions in priority order
        for ext in extensions + [None]:  # Try last with no extension
            # Prepare kwargs for this extension
            _kwargs = copy.deepcopy(kwargs)
            
            # Overwrite file_extension appropriately
            if ext is None:
                if "file_extension" in _kwargs:
                    del _kwargs["file_extension"]
            else:
                _kwargs.update({"file_extension": ext})
            converters_to_try = []
            
            if ext is not None and ext.lower() in self._extension_mapping:
                # Use mapped converters for known extensions
                converters_to_try = self._extension_mapping[ext.lower()]
            else:
                # For unknown extensions or None, try text converter first, then fallback to all converters
                converters_to_try = [self._text_converter]
                # If text converter fails, we'll fallback to the original approach below
            
            # Try the mapped converters first
            for converter in converters_to_try:
                try:
                    res = converter.convert(source_path, **_kwargs)
                    if res is not None:
                        # Normalize the content
                        res.text_content = "\n".join([line.rstrip() for line in re.split(r"\r?\n", res.text_content)])
                        res.text_content = re.sub(r"\n{3,}", "\n\n", res.text_content)
                        return res
                except Exception:
                    error_trace = ("\n\n" + traceback.format_exc()).strip()
            
            # If mapped converters failed and this is an unknown extension, fallback to trying all converters
            if ext is None or ext.lower() not in self._extension_mapping:
                for converter in self._page_converters:
                    # Skip converters we already tried
                    if converter in converters_to_try:
                        continue
                        
                    try:
                        res = converter.convert(source_path, **_kwargs)
                        if res is not None:
                            # Normalize the content
                            res.text_content = "\n".join([line.rstrip() for line in re.split(r"\r?\n", res.text_content)])
                            res.text_content = re.sub(r"\n{3,}", "\n\n", res.text_content)
                            return res
                    except Exception:
                        error_trace = ("\n\n" + traceback.format_exc()).strip()

        # If we got this far without success, report any exceptions
        if len(error_trace) > 0:
            raise FileConversionException(
                f"Could not convert '{source_path}' to Markdown. File type was recognized as {extensions}. While converting the file, the following error was encountered:\n\n{error_trace}"
            )

        # Nothing can handle it!
        raise UnsupportedFormatException(
            f"Could not parse the file:'{source_path}'. The formats {extensions} are not supported."
        )

    def _append_ext(self, extensions, ext):
        """Append a unique non-None, non-empty extension to a list of extensions."""
        if ext is None:
            return
        ext = ext.strip()
        if ext == "":
            return
        # if True
        if ext not in extensions:
            extensions.append(ext)

    def _guess_ext_magic(self, path):
        """Use puremagic (a Python implementation of libmagic) to guess a file's extension based on the first few bytes."""
        # Use puremagic to guess
        try:
            guesses = puremagic.magic_file(path)
            if len(guesses) > 0:
                ext = guesses[0].extension.strip()
                if len(ext) > 0:
                    return ext
        except FileNotFoundError:
            pass
        except IsADirectoryError:
            pass
        except PermissionError:
            pass
        return None

    def register_page_converter(self, converter: DocumentConverter) -> None:
        """Register a page text converter."""
        self._page_converters.insert(0, converter)
